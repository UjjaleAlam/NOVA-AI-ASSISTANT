from pptx import Presentation


def extract_text(path):
    """
    Extract all text from a PowerPoint presentation.
    """

    try:

        presentation = Presentation(path)

        text = []

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    content = shape.text.strip()

                    if content:

                        text.append(content)

        return "\n".join(text)

    except Exception as e:

        print(f"[PPTX Reader] {e}")

        return ""